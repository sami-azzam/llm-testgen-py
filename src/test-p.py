from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "LLM-Generated Test-Suite Evaluation for Defects4J"
subtitle.text = "Zero-Shot vs Chain-of-Thought on GPT-4o-mini\n& Comparison with EvoSuite\nMay 2025"

# Slide 2 – Methodology
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide.shapes.title, slide.shapes.placeholders[1].text_frame
title.text = "Methodology"
body.clear()
for t in [
    "Dataset: Defects4J (4 Java projects, buggy version “_1b”)",
    "Two prompt styles ➜ Zero-Shot (ZSL) vs Chain-of-Thought (CoT)",
    "LLM: GPT-4o-mini (≈ 4.1-nano)",
    "Pipeline:",
    "  1. Generate JUnit tests per class (≤82 kB) → generated-tests/<prompt>/…",
    "  2. Build & run with defects4j compile / test",
    "  3. Collect JaCoCo XML coverage & summarise (line, branch, method)",
]:
    p = body.add_paragraph()
    p.text = t
    p.level = 0 if t.startswith("Dataset") else 1 if t.startswith("Two") else 2 if t.startswith("  ") else 0
    p.font.size = Pt(18)

# Slide 3 – Results
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Coverage Results"
rows, cols = 6, 4
table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.5), Inches(9), Inches(3)).table

hdrs = ["Project", "Line % (ZSL / CoT)", "Branch % (ZSL / CoT)", "EvoSuite Line %"]
for i, h in enumerate(hdrs):
    cell = table.cell(0, i)
    cell.text = h

data = [
    ("JacksonCore_1b", "80.2 / 80.2", "—", "—"),
    ("Closure_1b", "98.2 / 98.2", "—", "—"),
    ("Csv_1b", "86.5 / 86.5", "—", "—"),
    ("JxPath_1b", "77.7 / 77.7", "—", "—"),
    ("AVERAGE", "85.65 (proj) / 90.21 (file)", "—", "78.48 (Gson)"),
]

for r, row in enumerate(data, 1):
    for c, val in enumerate(row):
        table.cell(r, c).text = str(val)

# Slide 4 – Tech Stack
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide.shapes.title, slide.shapes.placeholders[1].text_frame
title.text = "Technology Stack"
body.clear()
stack_items = [
    "TypeScript 5  •  Node.js 20  •  tsx runtime",
    "OpenAI SDK 4  •  GPT-4o-mini for generation",
    "Defects4J build wrappers (Ant/Maven) under JDK-11",
    "JaCoCo 0.8.11 for line/branch/method coverage",
    "xml2js, p-limit, minimist, chalk – CLI helpers",
]
for it in stack_items:
    p = body.add_paragraph()
    p.text = f"• {it}"
    p.font.size = Pt(20)

# Save file
file_path = "./LLM_Testgen_Slides.pptx"
prs.save(file_path)

file_path
